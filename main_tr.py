import comtypes.client
import fitz
import matplotlib.pyplot as plt
import os
import pandas as pd
import PyPDF2
import random
import string
import tabula
import tkinter as tk
import win32com.client
import webbrowser
from docx import Document
from docx.shared import Inches
from docx2pdf import convert
from matplotlib.backends.backend_pdf import PdfPages
from PIL import Image, ImageTk
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.pdfgen import canvas
from tkinter import filedialog, messagebox


# pdf birleştir
def merge_pdfs(pdf_paths, output_folder):
    pdf_writer = PyPDF2.PdfWriter()

    for pdf_path in pdf_paths:
        with open(pdf_path, "rb") as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                pdf_writer.add_page(page)

    random_filename = (
        "".join(random.choices(string.ascii_letters + string.digits, k=8)) + ".pdf"
    )
    output_path = os.path.join(output_folder, random_filename)

    with open(output_path, "wb") as output_file:
        pdf_writer.write(output_file)

    return output_path


# pdf ayır
def split_pdf(pdf_path, page_ranges, output_folder):
    page_ranges_list = page_ranges.split(",")
    pdf_reader = PyPDF2.PdfReader(pdf_path)
    total_pages = len(pdf_reader.pages)

    try:
        pdf_writer_list = []

        for page_range in page_ranges_list:
            if page_range.startswith("#"):
                step = int(page_range[1:])
                pdf_writer = PyPDF2.PdfWriter()

                for page_num in range(step - 1, total_pages, step):
                    page = pdf_reader.pages[page_num]
                    pdf_writer.add_page(page)

                pdf_writer_list.append(pdf_writer)

            else:
                page_range_parts = page_range.split("+")
                pdf_parts = []

                for part in page_range_parts:
                    part = part.replace(" ", "")
                    if "-" in part:
                        start, end = map(int, part.split("-"))

                        if not part.isdigit():
                            raise ValueError("0000: Invalid characters in page_ranges.")
                        elif start <= 0 or end > total_pages or start > end:
                            raise ValueError("0002: Invalid page range.")
                    else:
                        start = end = int(part)

                        if start <= 0 or end > total_pages:
                            raise ValueError("0001: Invalid page number.")

                    pdf_writer = PyPDF2.PdfWriter()

                    for page_num in range(total_pages):
                        page = pdf_reader.pages[page_num]

                        if start <= page_num + 1 <= end:
                            pdf_writer.add_page(page)

                    pdf_parts.append(pdf_writer)

                merged_writer = PyPDF2.PdfWriter()
                for part_writer in pdf_parts:
                    for page in part_writer.pages:
                        merged_writer.add_page(page)

                pdf_writer_list.append(merged_writer)

        for i, pdf_writer in enumerate(pdf_writer_list):
            split_path = f"{output_folder}/split_{i + 1}.pdf"
            with open(split_path, "wb") as split_file:
                pdf_writer.write(split_file)

            print(f"PDF split into {split_path}")

        return "Success"

    except ValueError as e:
        return str(e)


def convert_file_name(file_name, file_type, page_number=None):
    conversion_keyword = "_converted_to_"
    if conversion_keyword in file_name:
        file_name = file_name[
            : file_name.rfind(conversion_keyword)
        ]  # Remove last "converted_to" part

    return f"{file_name}{conversion_keyword}{file_type}{'_page' + str(page_number) if page_number else ''}"


# convert to from Word
def convert_to_opposite_format(input_file, output_folder):
    try:
        file_name, file_extension = os.path.splitext(os.path.basename(input_file))
        output_file = output_folder

        if file_extension.lower() == ".docx":
            output_path = os.path.join(
                output_file, convert_file_name(file_name, "PDF") + ".pdf"
            )
            convert(input_file, output_path)
            messagebox.showinfo(
                "Conversion Successful",
                f"File converted to PDF and saved to: {output_path}",
            )
            print(f"Conversion successful: {input_file} to {output_path}")

        elif file_extension.lower() == ".pdf":
            output_path = os.path.join(
                output_file, convert_file_name(file_name, "Word") + ".docx"
            )

            with fitz.open(input_file) as pdf_document:
                doc = Document()
                for page_number in range(pdf_document.page_count):
                    page = pdf_document[page_number]

                    # PDF'ten metni al ve DOCX'e ekle
                    text = page.get_text()
                    doc.add_paragraph(text)

                    # PDF'ten resimleri al ve DOCX'e ekle
                    images = page.get_images(full=True)
                    for img_index, img in enumerate(images):
                        img_index += 1
                        img_xref = img[0]
                        base_image = pdf_document.extract_image(img_xref)
                        image_bytes = base_image["image"]
                        image_filename = f"image_{page_number + 1}_{img_index}.png"
                        with open(image_filename, "wb") as image_file:
                            image_file.write(image_bytes)
                        doc.add_picture(image_filename, width=Inches(2.0))

            # Docx dosyasını kaydet
            doc.save(output_path)
            messagebox.showinfo(
                "Conversion Successful",
                f"File converted to Word and saved to: {output_path}",
            )
            print(f"Conversion successful: {input_file} to {output_path}")

        else:
            print(f"Unsupported file type: {file_extension}")

    except Exception as e:
        print(f"Conversion failed: {e}")


# Excel PDF dönüştür
def convert_excelpdf(input_file, output_folder):
    if input_file.lower().endswith((".xls", ".xlsx")):
        # Excel'i PDF'e dönüştür
        df = pd.read_excel(input_file)
        output_file = os.path.join(
            output_folder,
            f"{os.path.splitext(os.path.basename(input_file))[0]}_converted_to_PDF.pdf",
        )

        fig, ax = plt.subplots(figsize=(10, 6))
        ax.axis("tight")
        ax.axis("off")

        cell_colors = [
            ["#D3D3D3" if i == 0 else "white" for i in range(len(df.columns))]
        ] * len(df)

        table = ax.table(
            cellText=df.values,
            colLabels=df.columns,
            cellLoc="center",
            loc="center",
            cellColours=cell_colors,
        )

        table.auto_set_font_size(False)
        table.set_fontsize(10)

        for i, col in enumerate(df.columns):
            max_len = max(df[col].astype(str).apply(len).max(), len(col))
            col_width = max_len + 0.5
            table.auto_set_column_width([i])
            if table.get_celld()[(0, i)].get_width() < col_width:
                table.get_celld()[(0, i)].set_width(col_width)

        with PdfPages(output_file) as pdf:
            pdf.savefig(fig, bbox_inches="tight")

        plt.close()

        messagebox.showinfo(
            "Conversion Successful",
            f"File converted to PDF and saved to: {output_file}",
        )

    elif input_file.lower().endswith(".pdf"):
        # PDF'i Excel'e dönüştür
        dfs = tabula.read_pdf(input_file, pages="all")
        df = pd.concat(dfs, ignore_index=True)
        df = df.loc[:, ~df.columns.str.startswith("Unnamed: 6")]
        output_file = os.path.join(
            output_folder,
            f"{os.path.splitext(os.path.basename(input_file))[0]}_converted_to_Excel.xlsx",
        )
        df.to_excel(output_file, index=False)

        messagebox.showinfo(
            "Conversion Successful",
            f"File converted to Excel and saved to: {output_file}",
        )

    else:
        messagebox.showerror(
            "Unsupported File Type", f"Unsupported file type: {input_file}"
        )


# PDF'i PPTX'e dönüştür
def convert_file(input_file, output_folder):
    file_name_without_extension = os.path.splitext(os.path.basename(input_file))[0]
    output_file = os.path.join(
        output_folder, f"{file_name_without_extension}_converted"
    )

    if input_file.lower().endswith(".pdf"):
        # PDF'i PPTX'e dönüştür
        convert_pdf_to_pptx(input_file, output_file + ".pptx")
    elif input_file.lower().endswith((".ppt", ".pptx")):
        # PPTX'i PDF'e dönüştür
        convert_ppt_to_pdf(input_file, output_file + ".pdf")
    else:
        print("Unsupported file format. Please provide a PDF or PPTX file.")


def convert_pdf_to_pptx(input_pdf, output_pptx):
    try:
        presentation = Presentation()
        pdf_document = fitz.open(input_pdf)

        for page_number in range(pdf_document.page_count):
            page = pdf_document[page_number]
            text = page.get_text()

            slide = presentation.slides.add_slide(presentation.slide_layouts[5])

            text_box = slide.shapes.add_textbox(
                left=Inches(1), top=Inches(1), width=Inches(8), height=Inches(5)
            )
            text_frame = text_box.text_frame
            text_frame.text = text

        presentation.save(output_pptx)
        messagebox.showinfo(
            "Conversion Successful", f"{input_pdf} converted to PowerPoint."
        )
    except Exception as e:
        messagebox.showerror("Conversion Error", f"An error occurred: {e}")


def convert_ppt_to_pdf(input_pptx, output_pdf):
    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = True

        presentation = powerpoint.Presentations.Open(input_pptx)

        presentation.ExportAsFixedFormat(output_pdf, 32)  # 32 for PDF

        presentation.Close()
        powerpoint.Quit()
        messagebox.showinfo("Conversion Successful", f"{input_pptx} converted to PDF.")
    except Exception as e:
        messagebox.showerror("Conversion Error", f"An error occurred: {e}")


def ppt_pdf(input_path, output_path):
    if input_path.lower().endswith((".ppt", ".pptx")):
        # PPTX'i PDF'e dönüştür
        convert_ppt_to_pdf(input_path, output_path)
    elif input_path.lower().endswith(".pdf"):
        # PDF'i PPTX'e dönüştür
        convert_pdf_to_pptx(input_path, output_path)
    else:
        print("Unsupported file format. Please provide a PDF or PPTX file.")


# rtf-pdf dönüştür
def convert_to_pdf_or_rtf(input_file, output_folder):
    try:
        file_name, file_extension = os.path.splitext(os.path.basename(input_file))
        print(f"File extension: {file_extension}")

        output_file = os.path.join(output_folder, "")

        if file_extension.lower() == ".pdf":
            # PDF'i RTF'ye dönüştür
            with open(input_file, "rb") as pdf_file:
                doc = fitz.open(pdf_file)
                text_content = ""
                for page in doc:
                    text_content += page.get_text()

            # RTF dosyasını oluştur:
            docx_doc = Document()
            for paragraph in text_content.split("\n"):
                docx_doc.add_paragraph(paragraph)

            # RTF olarak kaydet:
            output_path = os.path.join(
                output_file, convert_file_name(file_name, "RTF") + ".rtf"
            )
            docx_doc.save(output_path)

            print(f"Conversion successful: {input_file} to {output_path}")

        elif file_extension.lower() == ".rtf":
            # RTF'yi PDF'ye dönüştür
            word = win32com.client.Dispatch("Word.Application")
            doc = word.Documents.Open(input_file)
            output_path = os.path.join(
                output_file, convert_file_name(file_name, "PDF") + ".pdf"
            )
            doc.SaveAs(output_path, FileFormat=17)  # 17 = wdFormatPDF
            word.Quit()

            print(f"Conversion successful: {input_file} to {output_path}")

        else:
            print(f"Unsupported file type: {file_extension}")

    except Exception as e:
        print(f"Conversion failed: {e}")


# txt-pdf dönüştür
def convert_to_pdf_or_txt():
    # Kullanıcı: bir dosya seçsin (PDF veya TXT)
    input_file = filedialog.askopenfilename(
        filetypes=[("PDF Files", "*.pdf"), ("Text Files", "*.txt")]
    )
    if input_file:
        # kullanıcı: bir hedef klasör seçsin
        output_folder = filedialog.askdirectory()
        if output_folder:
            try:
                file_name, file_extension = os.path.splitext(
                    os.path.basename(input_file)
                )
                output_path = os.path.join(
                    output_folder,
                    (
                        f"{file_name}.pdf"
                        if file_extension.lower() == ".txt"
                        else f"{file_name}.txt"
                    ),
                )

                if file_extension.lower() == ".pdf":
                    # PDF'i metin dosyasına dönüştür
                    with open(input_file, "rb") as pdf_file:
                        pdf_reader = PyPDF2.PdfReader(pdf_file)
                        text_content = ""
                        for page_number in range(len(pdf_reader.pages)):
                            text_content += pdf_reader.pages[page_number].extract_text()

                    # Metin dosyasını oluştur
                    with open(output_path, "w", encoding="utf-8") as text_file:
                        text_file.write(text_content)

                elif file_extension.lower() == ".txt":
                    # Metni PDF'ye dönüştür
                    doc = Document()
                    with open(input_file, "r", encoding="utf-8") as txt_file:
                        for line in txt_file:
                            doc.add_paragraph(line.strip())

                    doc.save(output_path)

                messagebox.showinfo(
                    "Conversion Successful",
                    f"File converted successfully and saved to: {output_path}",
                )

            except Exception as e:
                messagebox.showerror("Conversion Error", f"An error occurred: {e}")


# PNG-PDF dönüştür
def convert_file_name(file_name, file_type, page_number=None):
    conversion_keyword = "_converted_to_"
    base_file_name = file_name.split("_converted_to_")[
        0
    ]  # Extract base name (always look for conversion keyword)
    return f"{base_file_name}{conversion_keyword}{file_type}{'_page{page_number}' if page_number else ''}"


def convert_png_pdf(input_file, output_folder):
    try:
        file_extension = os.path.splitext(input_file)[1].lower()
        file_name = os.path.splitext(os.path.basename(input_file))[0]

        if file_extension == ".pdf":
            # PDF'yi PNG'ye dönüştür
            with fitz.open(input_file) as pdf_document:
                for page_number in range(pdf_document.page_count):
                    page = pdf_document[page_number]
                    image = page.get_pixmap()
                    img = Image.frombytes(
                        "RGB", [image.width, image.height], image.samples
                    )
                    output_path = os.path.join(
                        output_folder,
                        convert_file_name(
                            file_name,
                            "PNG",
                            None if pdf_document.page_count == 1 else page_number + 1,
                        )
                        + ".png",
                    )
                    img.save(output_path, "PNG")

            print(
                f"Conversion successful: {input_file} to PNG images in {output_folder}"
            )

        elif file_extension == ".png":
            # PNG'yi PDF'ye dönüştür
            with Image.open(input_file) as img:
                img = img.convert("RGB")
                output_pdf = (
                    os.path.join(output_folder, convert_file_name(file_name, "PDF"))
                    + ".pdf"
                )
                pdf = canvas.Canvas(output_pdf, pagesize=img.size)
                pdf.drawInlineImage(img, 0, 0, width=img.width, height=img.height)
                pdf.save()

            print(f"Conversion successful: {input_file} to {output_pdf}")

        else:
            print("Unsupported file type. Please provide either a PDF or PNG.")

    except Exception as e:
        print(f"Conversion failed: {e}")


# daha fazla dosya tipi eklenebilir veya tum dosya uzantıları
def select_files():
    files = filedialog.askopenfilenames(
        filetypes=[("PDF Files", "*.pdf"), ("Excel Files", "*.xls;*.xlsx")]
    )
    for file in files:
        listbox_files.insert(tk.END, file)


def clear_files():
    listbox_files.delete(0, tk.END)


def merge_files():
    pdf_paths = listbox_files.get(0, tk.END)
    output_folder = filedialog.askdirectory()
    if pdf_paths and output_folder:
        merged_pdf_path = merge_pdfs(pdf_paths, output_folder)
        messagebox.showinfo(
            "Merge PDFs", f"Merged PDF Files saved to: {merged_pdf_path}"
        )


# jpeg pdf dönüştür
def convert_to_pdf_or_jpeg(input_file, output_folder):
    try:
        file_name, file_extension = os.path.splitext(os.path.basename(input_file))
        output_file = os.path.join(output_folder, "")

        if file_extension.lower() == ".pdf":
            # PDF'yi JPEG'e dönüştür
            with fitz.open(input_file) as pdf_document:
                for page_number in range(pdf_document.page_count):
                    page = pdf_document[page_number]
                    image = page.get_pixmap()
                    img = Image.frombytes(
                        "RGB", [image.width, image.height], image.samples
                    )
                    output_path = os.path.join(
                        output_file,
                        convert_file_name(file_name, "JPEG", page_number + 1) + ".jpg",
                    )
                    img.save(output_path, "JPEG")
                    print(f"Conversion successful: {input_file} to {output_path}")

        elif file_extension.lower() in {".jpg", ".jpeg"}:
            # JPEG'yi PDF'e dönüştür
            output_path = os.path.join(
                output_file, convert_file_name(file_name, "PDF") + ".pdf"
            )
            pdf = canvas.Canvas(output_path)
            pdf.drawInlineImage(Image.open(input_file), 0, 0)
            pdf.save()
            print(f"Conversion successful: {input_file} to {output_path}")

    except Exception as e:
        print(f"Conversion failed: {e}")


def merge_files():
    pdf_paths = listbox_files.get(0, tk.END)
    output_folder = filedialog.askdirectory()
    if pdf_paths and output_folder:
        merged_pdf_path = merge_pdfs(pdf_paths, output_folder)
        # Dosya yolu - dosya adı - dosya uzantısı
        file_names = [os.path.basename(pdf_path) for pdf_path in pdf_paths]
        messagebox.showinfo(
            "Merge PDFs",
            f"Merged PDF Files ({', '.join(file_names)}) saved to: {merged_pdf_path}",
        )


def split_file():
    pdf_path = listbox_files.get(tk.ACTIVE)
    page_ranges = entry_page_ranges.get()
    output_folder = filedialog.askdirectory()
    if pdf_path and page_ranges and output_folder:
        result = split_pdf(pdf_path, page_ranges, output_folder)
        if result == "Success":
            # Dosya yolundan dosya adını al
            file_name = os.path.basename(pdf_path)
            messagebox.showinfo(
                "Split PDF",
                f"Split PDF files {file_name} saved to: " + output_folder,
            )
        else:
            messagebox.showerror("Split PDF Error", result)


# filedialogs for all conversions
def convert_pdf_to_word():
    input_file = filedialog.askopenfilename(
        filetypes=[("PDF Files", "*.pdf"), ("Word Files", "*.docx")]
    )
    if input_file:
        output_folder = filedialog.askdirectory()
        if output_folder:
            convert_to_opposite_format(input_file, output_folder)


def convert_pdf_to_excel():
    input_file = filedialog.askopenfilename(
        filetypes=[("PDF Files", "*.pdf"), ("Excel Files", "*.xls;*.xlsx")]
    )
    if input_file:
        output_folder = filedialog.askdirectory()
        if output_folder:
            convert_excelpdf(input_file, output_folder)


# pptx fonskiyonu için filediaolog
def convert_pdf_to_pptx_gui():
    input_file = filedialog.askopenfilename(
        filetypes=[("PDF Files", "*.pdf"), ("PPTX Files", "*.pptx")]
    )
    if input_file:
        output_folder = filedialog.askdirectory()
        if output_folder:
            convert_file(input_file, output_folder)


# rtf fonskiyonu için filediaolog
def convert_to_pdf_or_rtf():
    input_file = filedialog.askopenfilename(
        filetypes=[("PDF Files", "*.pdf"), ("RTF Files", "*.rtf")]
    )
    if input_file:
        output_folder = filedialog.askdirectory()
        if output_folder:
            convert_to_opposite_format(input_file, output_folder)


# jpeg pdf fonskiyonu için filediaolog
def convert_to_pdf_or_jpeg_gui():
    input_file = filedialog.askopenfilename(
        filetypes=[("PDF Files", "*.pdf"), ("JPEG Files", "*.jpg;*.jpeg")]
    )
    if input_file:
        output_folder = filedialog.askdirectory()
        if output_folder:
            convert_to_pdf_or_jpeg(input_file, output_folder)


# png pdf fonskiyonu için filediaolog
def convert_to_pdf_or_png():
    # kullanıcıdan PDF veya PNG dosyası seçmesini ise
    input_file = filedialog.askopenfilename(
        filetypes=[("PDF Files", "*.pdf"), ("PNG Files", "*.png")]
    )
    if input_file:
        # Kullanıcıdan cıktı klasörü seçmesini iste
        output_folder = filedialog.askdirectory()
        if output_folder:
            convert_png_pdf(input_file, output_folder)


# Menü fonksiyonları ve baglantılar
def open_link(event):
    webbrowser.open_new("https://torpdf.wordpress.com")


def open_readme():  # Tarayıcıya README.md dosyasını acar
    webbrowser.open_new("https://github.com/Kreytorn/TorPDF/blob/main/README.md")


def open_license():  # Tarayıcıya LICENSE dosyasını acar
    webbrowser.open_new("https://github.com/Kreytorn/TorPDF/blob/main/LICENSE")


def open_wiki():  # Tarayıcıda Wiki ve Hakkında sayfasını acar ve sayfasını acar
    webbrowser.open_new("https://github.com/Kreytorn/TorPDF/wiki")


def open_forum():  # Tarayıcıda Forum sayfasını acar
    webbrowser.open_new("https://github.com/Kreytorn/TorPDF/discussions")


# Arayüz oluşturma
root = tk.Tk()
root.title("TorPDF")
root.geometry("800x740")
root.configure(bg="black")

# Menu bar oluşturma
menubar = tk.Menu(root)
root.config(menu=menubar)

# Dosya menü
file_menu = tk.Menu(menubar, tearoff=0)
menubar.add_cascade(label="Dosya", menu=file_menu)
file_menu.add_command(label="Dosya Aç", command=select_files)
file_menu.add_separator()
file_menu.add_command(label="Çıkış", command=root.quit)

# Edit menü
edit_menu = tk.Menu(menubar, tearoff=0)
menubar.add_cascade(label="Düzenle", menu=edit_menu)
edit_menu.add_command(label="Birleştir", command=merge_files)
edit_menu.add_command(label="Ayır", command=split_file)
edit_menu.add_separator()
edit_menu.add_command(label="PDF<-->Word", command=merge_files)
edit_menu.add_command(label="PDF<-->Excel", command=merge_files)
edit_menu.add_command(label="PDF<-->PPT", command=merge_files)
edit_menu.add_command(label="PDF<-->RTF", command=merge_files)
edit_menu.add_command(label="PDF<-->RTF", command=merge_files)
edit_menu.add_separator()
edit_menu.add_command(label="PDF<-->JPEG", command=merge_files)
edit_menu.add_command(label="PDF<-->PNG", command=merge_files)
edit_menu.add_separator()
edit_menu.add_command(label="Temizle", command=clear_files)

# Araçlar menü
tools_menu = tk.Menu(menubar, tearoff=0)
menubar.add_cascade(label="Araçlar", menu=tools_menu)

# Yardım menü
help_menu = tk.Menu(menubar, tearoff=0)
menubar.add_cascade(label="Yardım", menu=help_menu)
help_menu.add_command(label="Hoşgeldiniz", command=open_readme)
help_menu.add_command(label="Kılavuz", command=open_readme)
help_menu.add_separator()
help_menu.add_command(label="Lisans", command=open_license)
help_menu.add_command(label="Soru&Cevap-Forum", command=open_forum)
help_menu.add_command(label="Hakkında", command=open_wiki)

frame_menu = tk.Frame(root, bg="black")
frame_menu.pack(side=tk.LEFT, fill=tk.Y)

# link ile logo ekleme
logo_image = Image.open("torpdf_logo.png")
logo_image = logo_image.resize((100, 100), Image.LANCZOS)
logo_photo = ImageTk.PhotoImage(logo_image)
logo_label = tk.Label(root, image=logo_photo, bg="black")
logo_label.place(x=20, rely=1.0, anchor="sw")

label_appname = tk.Label(
    frame_menu, text="TorPDF", bg="black", fg="white", font=("Helvetica", 20, "bold")
)
label_appname.pack(side=tk.TOP, pady=(0, 20))

label_version = tk.Label(frame_menu, text="Versiyon 1.0", bg="black", fg="white")
label_version.pack(side=tk.TOP)

# button içeren sol çerçeve
button_select = tk.Button(
    frame_menu,
    text="Birleştir-Ayır\nDosya seçin",
    bg="burlywood3",
    fg="black",
    command=select_files,
    width=13,
)
button_select.pack(pady=(50, 10), padx=20)

button_merge = tk.Button(
    frame_menu,
    text="Birleştir",
    bg="burlywood3",
    fg="black",
    command=merge_files,
    width=13,
)
button_merge.pack(pady=10, padx=20)

button_split = tk.Button(
    frame_menu, text="Ayır", bg="burlywood3", fg="black", command=split_file, width=13
)
button_split.pack(pady=10, padx=20)

# range - aralık
label_page_ranges = tk.Label(
    frame_menu,
    text="Ayırmak için sayfa aralık gir\n (e.g., 1-3, 5, 7-10, #2):",
    bg="black",
    fg="white",
)
label_page_ranges.pack()

entry_page_ranges = tk.Entry(frame_menu)
entry_page_ranges.pack(pady=(0, 20))

# PDF to Word (.docx)
button_pdftoword = tk.Button(
    frame_menu,
    text="PDF <--> Word",
    bg="cadetblue1",
    fg="black",
    command=convert_pdf_to_word,
    width=13,
)
button_pdftoword.pack(pady=10, padx=20)

# PDF to Excel (.xlsx)
button_pdftoexcel = tk.Button(
    frame_menu,
    text="PDF <--> Excel",
    bg="cadetblue1",
    fg="black",
    command=convert_pdf_to_excel,
    width=13,
)
button_pdftoexcel.pack(pady=10, padx=20)

# PDF to PowerPoint (.pptx)
button_pdftoppt = tk.Button(
    frame_menu,
    text="PDF <--> Pptx",
    bg="cadetblue1",
    fg="black",
    command=convert_pdf_to_pptx_gui,
    width=13,
)
button_pdftoppt.pack(pady=10, padx=20)

# PDF to Rtf (.rtf)
button_pdftortf = tk.Button(
    frame_menu,
    text="PDF <--> Rtf",
    bg="cadetblue1",
    fg="black",
    command=convert_to_pdf_or_rtf,
    width=13,
)
button_pdftortf.pack(pady=10, padx=20)

# PDF to Text (.txt)
button_pdftotext = tk.Button(
    frame_menu,
    text="PDF <--> Text",
    bg="cadetblue1",
    fg="black",
    command=convert_to_pdf_or_txt,
    width=13,
)
button_pdftotext.pack(pady=10, padx=20)

# PDF to Jpeg(jpeg)
button_pdftojpg = tk.Button(
    frame_menu,
    text="PDF <--> Jpeg",
    bg="cadetblue1",
    fg="black",
    command=lambda: convert_to_pdf_or_jpeg_gui(),
    width=13,
)
button_pdftojpg.pack(pady=10, padx=20)

# PDF to PNG (png, pdf)
button_pdftopng = tk.Button(
    frame_menu,
    text="PDF <--> Png",
    bg="cadetblue1",
    fg="black",
    command=convert_to_pdf_or_png,
    width=13,
)
button_pdftopng.pack(pady=10, padx=20)

# Nasıl kullanılır içeriği ile sağ çerçeve
frame_files = tk.Frame(root, bg="burlywood3")
frame_files.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

button_clear = tk.Button(frame_files, text="Panoyu Temizle", command=clear_files)
button_clear.pack(pady=15)

listbox_files = tk.Listbox(frame_files, width=50, height=20)
listbox_files.pack(fill=tk.BOTH, expand=True, padx=20, pady=(0, 20))

link_label = tk.Label(
    frame_files, text="Web sitemizi ziyaret edin!", fg="blue", cursor="hand2"
)
link_label.pack(pady=(0, 20))
link_label.bind("<Button-1>", open_link)

# Açıklamalar için metin widget
text_explanations = tk.Text(
    root,
    bg="black",
    fg="white",
    wrap="word",
    width=10,
    font=("Calibri", 11),
)
text_explanations.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=20, pady=20)

# Açıklamaları metin widget'a ekleme
text_explanations.insert(
    tk.END,
    "TORPDF NASIL KULLANILIR?\n\nTORPDF bir PDF yönetim aracıdır.\nAşağıda TORPDF'nin özelliklerini etkili bir şekilde kullanmanız için kısa bir rehber bulunmaktadır.\n\n",
)
text_explanations.insert(
    tk.END,
    "Dosyaları Seçin: \nDeğiştirmek istediğiniz PDF dosyalarını seçerek başlayın. Dosyaları birleştirmek veya bölmek/ayırmak için, bu ilk adım dosya listesini doldurur. Dosyalar seçildikten sonra Birleştirme veya Bölme işlemlerine devam edin.\n\n",
)
text_explanations.insert(
    tk.END,
    "Birleştirme: \nBirden fazla PDF'yi tek bir dosyada birleştirin. Birleştirmeyi, sayfa çıkarmak için aralık işlevini kullanın.\n\n",
)
text_explanations.insert(
    tk.END,
    "Bölme-Ayırma: \nBir PDF'yi birden fazla dosyaya bölün/ayırın. Bölme/Ayırma işlemini ayarlamak için aralık işlevini kullanın, sayfaları çıkarın.\n\n",
)
text_explanations.insert(
    tk.END,
    "PDF'ye Dönüştür:\nDosyalarınızı şu dosya biçimlerine dönüştürün: (.docx,. xlsx, .ppxt, .txt, .rtf, .odf, .odt, .ods, .jpeg, .png dosyalarını PDF'e).\n\n",
)
text_explanations.insert(
    tk.END,
    "PDF'den Dönüştür:\nDosyalarınızı şu dosya biçimlerinden dönüştürün: (.docx,. xlsx, .pptx, .txt, .rtf, .odf, .odt, .ods, jpeg, .png dosyalarını PDF'e).\n\n",
)

text_explanations.insert(
    tk.END,
    "WIKI SAYFALARI & BELGELER & İLETİŞİM:\n\nDaha fazla bilgi için: \nhttps://github.com/Kreytorn/TorPDF/wiki \n\nBelgeler için:\nhttps://github.com/Kreytorn/TorPDF.\n\nForumlar ve sorular:\nhttps://github.com/Kreytorn/TorPDF/discussions\n\nBize ulaşın:\nhttps://github.com/Kreytorn/TorPDF/issues",
)

root.mainloop()
